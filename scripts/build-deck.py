#!/usr/bin/env python3
"""Build Uplift Summer 2026 Onboarding Deck as a PPTX for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Dimensions: 16:9 widescreen ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Brand colors ─────────────────────────────────────────────────────────────
PURPLE_DARK   = RGBColor(0x2D, 0x1F, 0x6E)   # #2d1f6e deep navy-purple (bg)
PURPLE_MID    = RGBColor(0x5C, 0x4E, 0xB5)   # #5c4eb5 primary purple
PURPLE_LIGHT  = RGBColor(0xA0, 0x90, 0xE0)   # #a090e0 light purple accent
LAVENDER      = RGBColor(0xF0, 0xEC, 0xFF)   # #f0ecff soft lavender
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xF8, 0xF6, 0xFF)
GOLD          = RGBColor(0xF5, 0xC5, 0x42)   # warm gold for accents
GREEN         = RGBColor(0x22, 0xA3, 0x66)   # milestone green
AMBER         = RGBColor(0xD9, 0x77, 0x06)   # amber
SLATE         = RGBColor(0xB0, 0xA8, 0xCC)   # muted text

# ── Helper: blank slide (no layout placeholders) ─────────────────────────────
def blank_slide():
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)

# ── Helper: fill slide background ────────────────────────────────────────────
def fill_bg(slide, color):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Logo paths ───────────────────────────────────────────────────────────────
LOGO_UPLIFT_WHITE = "/Users/techunited/Desktop/uplift-app/public/uplift-logo-white.png"
LOGO_UPLIFT_DARK  = "/Users/techunited/Desktop/uplift-app/public/uplift-logo.png"
LOGO_TECHUNITED   = "/Users/techunited/Desktop/uplift-app/public/techunited-logo.png"

# ── Helper: add logos to a slide ─────────────────────────────────────────────
def add_logos(slide, dark=True, uplift_width=Inches(1.1)):
    uplift_logo = LOGO_UPLIFT_WHITE if dark else LOGO_UPLIFT_DARK
    slide.shapes.add_picture(uplift_logo,
                             Inches(0.35), Inches(0.12),
                             width=uplift_width)
    slide.shapes.add_picture(LOGO_TECHUNITED,
                             W - Inches(1.65), H - Inches(0.52),
                             width=Inches(1.35))

# ── Helper: add rectangle ────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ── Helper: add text box ─────────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox

# ── Helper: multi-line text (list of (text, size, bold, color, align) tuples)
def add_multiline(slide, lines, x, y, w, h, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align, italic) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txBox

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True, uplift_width=Inches(1.4))

# Large gradient-feel left accent bar
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)

# Gold accent strip top
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.12), GOLD)

# Subtle right rectangle decoration
add_rect(s, W - Inches(3.2), Inches(2.2), Inches(3.2), Inches(3.1), RGBColor(0x3D, 0x2F, 0x8E))

# TechUnited label top-right
add_text(s, "TECHUNITED:NJ", W - Inches(3.0), Inches(0.3), Inches(2.8), Inches(0.5),
         size=11, bold=True, color=PURPLE_LIGHT, align=PP_ALIGN.RIGHT)

# Program title
add_text(s, "UPLIFT", Inches(0.55), Inches(1.4), Inches(7), Inches(1.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "SUMMER 2026", Inches(0.55), Inches(2.7), Inches(7), Inches(0.9),
         size=36, bold=False, color=PURPLE_LIGHT, align=PP_ALIGN.LEFT)

# Gold divider line
add_rect(s, Inches(0.55), Inches(3.6), Inches(4.5), Inches(0.05), GOLD)

# Subtitle
add_text(s, "Founder Onboarding", Inches(0.55), Inches(3.75), Inches(7), Inches(0.7),
         size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Welcome — and let's get started.", Inches(0.55), Inches(4.35), Inches(7), Inches(0.5),
         size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=True)

# Date pill
add_rect(s, Inches(0.55), Inches(5.4), Inches(2.8), Inches(0.5), PURPLE_MID)
add_text(s, "June 1 – August 4, 2026", Inches(0.55), Inches(5.4), Inches(2.8), Inches(0.5),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_WHITE)
add_logos(s, dark=False)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Today's Agenda", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)

items = [
    ("01", "Welcome & introductions"),
    ("02", "What is Uplift? — program overview & timeline"),
    ("03", "Meet your cohort"),
    ("04", "Program requirements — what's expected of you"),
    ("05", "The portal — how to find everything you need"),
    ("06", "Icebreakers & introductions"),
    ("07", "Q&A"),
]

for i, (num, label) in enumerate(items):
    y = Inches(1.3) + i * Inches(0.76)
    add_rect(s, Inches(0.5), y + Inches(0.08), Inches(0.48), Inches(0.48), PURPLE_MID)
    add_text(s, num, Inches(0.5), y + Inches(0.08), Inches(0.48), Inches(0.48),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, Inches(1.1), y, Inches(10), Inches(0.65),
             size=18, bold=False, color=PURPLE_DARK, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS UPLIFT
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "What Is Uplift?", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s,
    "A 9-week accelerator program for early-stage founders based in New Jersey — built around mentorship, "
    "peer community, and the focused work of moving your company forward.",
    Inches(0.5), Inches(1.1), Inches(8.5), Inches(1.1),
    size=18, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

pillars = [
    ("🧑‍🏫", "Dedicated Mentor", "1-on-1 sessions with an experienced operator or investor who knows your industry."),
    ("👥", "Peer Network", "A tight cohort of founders at a similar stage — real relationships, not just networking."),
    ("🎯", "Structured Progress", "Weekly reflection prompts, milestones, and a portal to keep you accountable."),
]

for i, (icon, title, desc) in enumerate(pillars):
    x = Inches(0.45) + i * Inches(4.22)
    add_rect(s, x, Inches(2.5), Inches(3.9), Inches(3.6), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, icon, x + Inches(0.2), Inches(2.65), Inches(0.8), Inches(0.7), size=28, color=WHITE)
    add_text(s, title, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(0.55),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, desc, x + Inches(0.2), Inches(3.85), Inches(3.55), Inches(1.7),
             size=13, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROGRAM TIMELINE
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_WHITE)
add_logos(s, dark=False)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "9 Weeks at a Glance", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)
add_text(s, "June 1 – August 4, 2026", Inches(0.5), Inches(1.0), Inches(5), Inches(0.4),
         size=15, bold=False, color=PURPLE_MID, align=PP_ALIGN.LEFT)

weeks = [
    ("Week 1",  "Jun 1–6",   "Onboarding",          "Orientation session · Portal access · Mentor reveal"),
    ("Weeks 2–4","Jun 8–20","Building Momentum",    "First mentor meeting · Educational sessions begin"),
    ("Week 4",  "Jun 23",   "Midpoint Meetup ★",   "Required in-person event · Halfway celebration"),
    ("Weeks 5–7","Jun 25–",  "Deep Work",            "Mentor sessions continue · Reflection prompts"),
    ("Week 8",  "Jul 19–25","Final Stretch",         "End of Program Survey due"),
    ("Week 9",  "Aug 4",    "Summit & Graduation ★","Final showcase · Networking · Certificate ceremony"),
]

col_w = Inches(2.1)
for i, (week, dates, title, detail) in enumerate(weeks):
    x = Inches(0.4) + i * col_w
    is_special = "★" in title
    bg_color = PURPLE_MID if is_special else RGBColor(0xED, 0xE9, 0xF8)
    txt_color = WHITE if is_special else PURPLE_DARK
    sub_color = RGBColor(0xD8, 0xD0, 0xFF) if is_special else SLATE

    add_rect(s, x, Inches(1.55), col_w - Inches(0.12), Inches(4.9), bg_color)
    add_text(s, week,  x + Inches(0.12), Inches(1.65), col_w - Inches(0.24), Inches(0.4),
             size=11, bold=True, color=sub_color, align=PP_ALIGN.LEFT)
    add_text(s, dates, x + Inches(0.12), Inches(2.0),  col_w - Inches(0.24), Inches(0.35),
             size=10, bold=False, color=sub_color, align=PP_ALIGN.LEFT)
    add_rect(s, x + Inches(0.12), Inches(2.3), Inches(0.3), Inches(0.03),
             WHITE if is_special else PURPLE_LIGHT)
    add_text(s, title, x + Inches(0.12), Inches(2.4), col_w - Inches(0.24), Inches(0.6),
             size=13, bold=True, color=txt_color, align=PP_ALIGN.LEFT)
    add_text(s, detail, x + Inches(0.12), Inches(3.0), col_w - Inches(0.24), Inches(2.3),
             size=11, bold=False, color=sub_color, align=PP_ALIGN.LEFT)

add_text(s, "★ Required in-person events", Inches(0.5), Inches(6.9), Inches(6), Inches(0.4),
         size=11, bold=False, color=SLATE, align=PP_ALIGN.LEFT, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE 5 COHORTS (title)
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Meet the Cohorts", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Each cohort is named for a trailblazer with roots in New Jersey.",
         Inches(0.5), Inches(1.1), Inches(11), Inches(0.55),
         size=17, bold=False, color=PURPLE_LIGHT, align=PP_ALIGN.LEFT)

cohorts = [
    ("1", "Edison",   "Go-to-market · Customer acquisition · Ship & iterate"),
    ("2", "Hopper",   "Enterprise SaaS · B2B · MVP-stage builders"),
    ("3", "Bardeen",  "AI / Data / ML · Technical founders"),
    ("4", "Lawrence", "Fundraising · Major raises · Pivots"),
    ("5", "Morrison", "Clarity · Narrative · Women & founders of color"),
]

for i, (num, name, tagline) in enumerate(cohorts):
    x = Inches(0.45) + i * Inches(2.55)
    add_rect(s, x, Inches(1.9), Inches(2.35), Inches(4.3), RGBColor(0x3D, 0x2F, 0x8E))
    add_rect(s, x, Inches(1.9), Inches(2.35), Inches(0.5), PURPLE_MID)
    add_text(s, num, x + Inches(0.12), Inches(1.92), Inches(0.4), Inches(0.45),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, name, x + Inches(0.15), Inches(2.5), Inches(2.1), Inches(0.55),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s, x + Inches(0.15), Inches(3.05), Inches(0.4), Inches(0.04), GOLD)
    add_text(s, tagline, x + Inches(0.15), Inches(3.15), Inches(2.1), Inches(2.7),
             size=12, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — COHORT DETAILS (one per cohort, combined on two slides)
# ═════════════════════════════════════════════════════════════════════════════
cohort_details = [
    {
        "num": "Cohort 1", "name": "Edison",
        "namesake": "Thomas Edison",
        "bio": "Born in Ohio but built his legacy in New Jersey. His Menlo Park lab was the original 'invention factory.' Edison held 1,093 U.S. patents.",
        "why": "The operator-heavy, customer-chasing cohort. Edison's genius wasn't just invention — it was relentless iteration toward products real customers would buy.",
        "theme": "Go-to-Market · Customer Acquisition",
        "color": RGBColor(0x5C, 0x4E, 0xB5),
    },
    {
        "num": "Cohort 2", "name": "Hopper",
        "namesake": "Grace Hopper",
        "bio": "Computer-science pioneer who worked at Sperry Rand in Bloomfield, NJ. Invented the first compiler and was a driving force behind COBOL.",
        "why": "The builders' cohort. Heavy on Enterprise SaaS and B2B. Hopper believed in building the tools that let everyone else build.",
        "theme": "Enterprise SaaS · MVP-Stage Builders",
        "color": RGBColor(0x1A, 0x7A, 0xB8),
    },
    {
        "num": "Cohort 3", "name": "Bardeen",
        "namesake": "John Bardeen",
        "bio": "Physicist at Bell Labs in Murray Hill, NJ — co-invented the transistor in 1947. Two-time Nobel Prize winner in Physics.",
        "why": "The most technical cohort — AI, Data, and ML founders. Bell Labs in NJ is literally where modern computing began.",
        "theme": "AI · Data · ML",
        "color": RGBColor(0x17, 0x8A, 0x6E),
    },
    {
        "num": "Cohort 4", "name": "Lawrence",
        "namesake": "Jacob Lawrence",
        "bio": "Painter born in Atlantic City, NJ. Best known for the Migration Series — 60 panels chronicling communities in moments of transformation.",
        "why": "The inflection-point cohort. Highest concentration of founders focused on major raises, launches, and pivots. Lawrence's art was about transition and becoming.",
        "theme": "Fundraising · Pivots · Major Launches",
        "color": RGBColor(0xA0, 0x5C, 0x1A),
    },
    {
        "num": "Cohort 5", "name": "Morrison",
        "namesake": "Toni Morrison",
        "bio": "Nobel laureate in Literature who lived and wrote in Princeton, NJ for decades, teaching at Princeton University.",
        "why": "The clarity-seeking cohort. Founders focused on clarifying priorities and refining narrative. Strongest representation of women and founders of color.",
        "theme": "Clarity · Narrative · Voice",
        "color": RGBColor(0x8B, 0x35, 0x9A),
    },
]

for cd in cohort_details:
    s = blank_slide()
    fill_bg(s, PURPLE_DARK)
    add_logos(s, dark=True)

    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, cd["color"])
    add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

    # Cohort label chip
    add_rect(s, Inches(0.5), Inches(0.35), Inches(1.3), Inches(0.38), cd["color"])
    add_text(s, cd["num"], Inches(0.5), Inches(0.35), Inches(1.3), Inches(0.38),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Name
    add_text(s, cd["name"], Inches(2.0), Inches(0.25), Inches(8), Inches(0.8),
             size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Theme pill
    add_rect(s, Inches(2.0), Inches(1.05), Inches(4.5), Inches(0.38), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, cd["theme"], Inches(2.0), Inches(1.05), Inches(4.5), Inches(0.38),
             size=12, bold=True, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)

    # Left column: bio
    add_rect(s, Inches(0.5), Inches(1.65), Inches(5.8), Inches(4.5), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, "About " + cd["namesake"], Inches(0.7), Inches(1.8), Inches(5.4), Inches(0.5),
             size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(s, cd["bio"], Inches(0.7), Inches(2.3), Inches(5.4), Inches(3.5),
             size=14, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

    # Right column: why this cohort
    add_rect(s, Inches(6.6), Inches(1.65), Inches(6.45), Inches(4.5), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, "Why this cohort?", Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.5),
             size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(s, cd["why"], Inches(6.8), Inches(2.3), Inches(6.0), Inches(3.5),
             size=14, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PROGRAM REQUIREMENTS
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_WHITE)
add_logos(s, dark=False)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Program Requirements", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)
add_text(s, "8 things to complete by August 4 to graduate from Uplift.",
         Inches(0.5), Inches(1.0), Inches(10), Inches(0.45),
         size=15, bold=False, color=PURPLE_MID, align=PP_ALIGN.LEFT)

reqs = [
    ("1", "Confirm participation",      "By Jun 3"),
    ("2", "Attend onboarding session",  "By Jun 7"),
    ("3", "3 educational sessions",     "By Aug 4"),
    ("4", "3 mentorship sessions",      "Jun 13 · Jul 4 · Jul 18"),
    ("5", "Midpoint Meetup",            "Jun 23 ★"),
    ("6", "End of Program Survey",      "By Jul 25"),
    ("7", "Summit & Graduation",        "Aug 4 ★"),
    ("8", "Receive certificate",        "Upon completion"),
]

col = 0
for i, (num, label, due) in enumerate(reqs):
    col = i % 2
    row = i // 2
    x = Inches(0.45) + col * Inches(6.35)
    y = Inches(1.6) + row * Inches(1.2)

    is_special = "★" in due
    bg = LAVENDER if not is_special else RGBColor(0xFE, 0xF3, 0xC7)
    due_color = AMBER if is_special else PURPLE_MID

    add_rect(s, x, y, Inches(6.05), Inches(1.05), bg)
    # Number circle
    add_rect(s, x + Inches(0.12), y + Inches(0.28), Inches(0.45), Inches(0.45), PURPLE_MID)
    add_text(s, num, x + Inches(0.12), y + Inches(0.28), Inches(0.45), Inches(0.45),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Label
    add_text(s, label, x + Inches(0.7), y + Inches(0.18), Inches(3.7), Inches(0.65),
             size=14, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)
    # Due date
    add_text(s, due, x + Inches(4.5), y + Inches(0.22), Inches(1.45), Inches(0.55),
             size=12, bold=True, color=due_color, align=PP_ALIGN.RIGHT)

add_text(s, "Track your progress in real-time on your personal portal.", Inches(0.5), Inches(6.85), Inches(10), Inches(0.45),
         size=12, bold=False, color=SLATE, italic=True, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — REQUIREMENTS DETAIL
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "A Few Things to Know", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

notes = [
    ("📅", "Sessions happen every Tuesday",
     "Educational sessions are held weekly on Tuesdays. You need 3 total — attend any session across any cohort. Schedules are unpredictable; show up when you can."),
    ("🎙️", "Granola is required for mentor sessions",
     "Download Granola (granola.ai) — a free AI note-taker. Run it during every mentorship session. Your key takeaways are part of the submission."),
    ("💻", "Submit through your portal",
     "After each mentor session, submit your meeting report through your personal portal. This is how we track your progress and keep you on track to graduate."),
]

for i, (icon, title, detail) in enumerate(notes):
    x = Inches(0.45) + i * Inches(4.22)
    add_rect(s, x, Inches(1.5), Inches(3.95), Inches(4.6), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, icon, x + Inches(0.2), Inches(1.65), Inches(0.8), Inches(0.7), size=28, color=WHITE)
    add_text(s, title, x + Inches(0.2), Inches(2.35), Inches(3.55), Inches(0.6),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s, x + Inches(0.2), Inches(2.95), Inches(0.5), Inches(0.04), GOLD)
    add_text(s, detail, x + Inches(0.2), Inches(3.1), Inches(3.55), Inches(2.7),
             size=13, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THE PORTAL
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_WHITE)
add_logos(s, dark=False)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Your Personal Portal", Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
         size=34, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)
add_text(s, "uplift2026.vercel.app / [your-name]",
         Inches(0.5), Inches(1.05), Inches(7), Inches(0.45),
         size=16, bold=False, color=PURPLE_MID, align=PP_ALIGN.LEFT)

tabs = [
    ("🗓️", "My Journey", "Week-by-week content, reflection prompts, and session links."),
    ("✅", "Milestones", "Live progress tracker — see what's done and what's coming up."),
    ("💡", "My Goals &\nReflections", "Everything you've written in one place — your personal journal."),
    ("👩‍💼", "Logged Mentor\nSessions", "Submit and track your mentor meetings here."),
    ("🎓", "Educational\nSessions", "Log attendance for your 3 required educational sessions."),
    ("🗺️", "Program\nRoadmap", "Full program schedule and key dates at a glance."),
    ("👥", "Cohort\nDirectory", "Get to know your fellow founders."),
    ("📚", "Resources", "Guides, templates, tools, and links — all in one place."),
]

col_count = 4
for i, (icon, name, desc) in enumerate(tabs):
    col = i % col_count
    row = i // col_count
    x = Inches(0.45) + col * Inches(3.2)
    y = Inches(1.65) + row * Inches(2.35)
    add_rect(s, x, y, Inches(3.0), Inches(2.15), LAVENDER)
    add_text(s, icon, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.55), size=22, color=PURPLE_DARK)
    add_text(s, name, x + Inches(0.15), y + Inches(0.7), Inches(2.7), Inches(0.65),
             size=13, bold=True, color=PURPLE_DARK, align=PP_ALIGN.LEFT)
    add_text(s, desc, x + Inches(0.15), y + Inches(1.3), Inches(2.7), Inches(0.75),
             size=11, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE — TECHUNITED: WHO WE ARE
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "About TechUnited:NJ", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s,
    "The largest tech community in New Jersey — connecting founders, operators, investors, "
    "and builders across the state.",
    Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.75),
    size=17, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

tu_points = [
    ("🤝", "Community First",
     "TechUnited is home to thousands of founders, operators, and innovators across New Jersey. "
     "We exist to connect, support, and amplify the people building the future here."),
    ("🚀", "We Run Programs",
     "Uplift is one of several programs TechUnited runs to support founders at every stage — "
     "from idea to growth. You're part of something bigger than this cohort."),
    ("🌐", "The NJ Ecosystem",
     "New Jersey is one of the most underrated startup ecosystems in the country. "
     "Dense talent, proximity to NYC capital, world-class universities — and TechUnited is the connective tissue."),
]

for i, (icon, title, body) in enumerate(tu_points):
    x = Inches(0.45) + i * Inches(4.22)
    add_rect(s, x, Inches(2.1), Inches(3.95), Inches(4.2), RGBColor(0x3D, 0x2F, 0x8E))
    add_text(s, icon, x + Inches(0.2), Inches(2.25), Inches(0.8), Inches(0.7), size=26, color=WHITE)
    add_text(s, title, x + Inches(0.2), Inches(2.95), Inches(3.55), Inches(0.55),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s, x + Inches(0.2), Inches(3.5), Inches(0.45), Inches(0.04), GOLD)
    add_text(s, body, x + Inches(0.2), Inches(3.65), Inches(3.55), Inches(2.4),
             size=12, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

add_text(s, "techunited.co", Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
         size=13, bold=True, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE — AI DEMO NIGHTS
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, GOLD)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

# "NEW" chip top-right
add_rect(s, W - Inches(1.8), Inches(0.25), Inches(1.3), Inches(0.38), GREEN)
add_text(s, "BONUS", W - Inches(1.8), Inches(0.25), Inches(1.3), Inches(0.38),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "AI Demo Nights", Inches(0.5), Inches(0.3), Inches(11), Inches(0.9),
         size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Hosted by TechUnited:NJ",
         Inches(0.5), Inches(1.15), Inches(7), Inches(0.45),
         size=16, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

add_text(s,
    "TechUnited runs regular AI Demo Nights — community events where founders, developers, "
    "and enthusiasts come together to demo what they're building with AI. "
    "As Uplift founders, you're encouraged to attend, present, and connect.",
    Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.0),
    size=16, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

demo_points = [
    ("🎤", "Demo your product",
     "AI Demo Nights are a low-stakes, high-energy place to put your product in front of a real audience. "
     "Great practice before investor pitches."),
    ("👀", "See what others are building",
     "Stay sharp on what's being built in the AI space right now — "
     "in your backyard, by people who might become customers, partners, or co-founders."),
    ("🤝", "Grow your network",
     "The room is full of technical talent, operators, and investors who care about AI. "
     "These are your future hires, advisors, and advocates."),
]

for i, (icon, title, body) in enumerate(demo_points):
    x = Inches(0.45) + i * Inches(4.22)
    add_rect(s, x, Inches(2.85), Inches(3.95), Inches(3.95), RGBColor(0x2A, 0x1F, 0x60))
    add_rect(s, x, Inches(2.85), Inches(3.95), Inches(0.06), GOLD)
    add_text(s, icon, x + Inches(0.2), Inches(3.0), Inches(0.8), Inches(0.65), size=26, color=WHITE)
    add_text(s, title, x + Inches(0.2), Inches(3.65), Inches(3.55), Inches(0.55),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, body, x + Inches(0.2), Inches(4.2), Inches(3.55), Inches(2.4),
             size=12, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

add_rect(s, Inches(2.8), Inches(6.55), Inches(7.75), Inches(0.72), RGBColor(0x2A, 0x1F, 0x60))
add_text(s, "Apply to demo  →  form.typeform.com/to/voMHQcT0",
         Inches(2.8), Inches(6.55), Inches(7.75), Inches(0.72),
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ICEBREAKER: INTRO FORMAT
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, GOLD)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Let's Meet Each Other", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Take 60 seconds. Share these four things:",
         Inches(0.5), Inches(1.05), Inches(10), Inches(0.5),
         size=18, bold=False, color=PURPLE_LIGHT, align=PP_ALIGN.LEFT)

intros = [
    ("👋", "Your name", "First name is fine."),
    ("🏢", "Your company", "What you're building — one sentence."),
    ("🔧", "What you do", "Are you technical? Sales-focused? Solo founder?"),
    ("🎯", "One goal for the summer", "What does success look like for you in 9 weeks?"),
]

for i, (icon, title, prompt) in enumerate(intros):
    y = Inches(1.75) + i * Inches(1.2)
    add_rect(s, Inches(0.5), y, Inches(12.35), Inches(1.05), RGBColor(0x3D, 0x2F, 0x8E))
    add_rect(s, Inches(0.5), y, Inches(0.7), Inches(1.05), PURPLE_MID)
    add_text(s, icon, Inches(0.5), y + Inches(0.15), Inches(0.7), Inches(0.7), size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.35), y + Inches(0.12), Inches(3.5), Inches(0.5),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, prompt, Inches(1.35), y + Inches(0.57), Inches(11.0), Inches(0.4),
             size=13, bold=False, color=SLATE, align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ICEBREAKER: DISCUSSION QUESTIONS
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, GOLD)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Discussion Questions", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Pick one. Go around the room. No wrong answers.",
         Inches(0.5), Inches(1.05), Inches(10), Inches(0.45),
         size=16, bold=False, color=PURPLE_LIGHT, align=PP_ALIGN.LEFT, italic=True)

questions = [
    "What's the one thing you most want to get out of this program — mentor access, peer connections, or accountability?",
    "What's the hardest part of building your company right now — and what would change if you solved it?",
    "If a mentor could wave a magic wand and fix one thing about your company this summer, what would it be?",
    "What's something you've tried that totally didn't work — and what did you learn from it?",
    "What does success look like for you on August 4th, at graduation?",
]

for i, q in enumerate(questions):
    y = Inches(1.65) + i * Inches(1.05)
    add_rect(s, Inches(0.5), y, Inches(12.35), Inches(0.92), RGBColor(0x3D, 0x2F, 0x8E))
    add_rect(s, Inches(0.5), y, Inches(0.42), Inches(0.92), GOLD)
    add_text(s, str(i+1), Inches(0.5), y + Inches(0.2), Inches(0.42), Inches(0.5),
             size=14, bold=True, color=PURPLE_DARK, align=PP_ALIGN.CENTER)
    add_text(s, q, Inches(1.05), y + Inches(0.14), Inches(11.6), Inches(0.65),
             size=14, bold=False, color=RGBColor(0xD8, 0xD0, 0xFF), align=PP_ALIGN.LEFT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING / Q&A
# ═════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, PURPLE_DARK)
add_logos(s, dark=True)

add_rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE_MID)
add_rect(s, Inches(0.18), Inches(0), W - Inches(0.18), Inches(0.08), GOLD)

add_text(s, "Questions?", Inches(0.5), Inches(1.8), Inches(12), Inches(1.5),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(4.5), Inches(3.35), Inches(4.35), Inches(0.06), GOLD)
add_text(s, "uplift@techunited.co", Inches(0.5), Inches(3.55), Inches(12), Inches(0.6),
         size=18, bold=False, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "We're glad you're here. Let's make this summer count.",
         Inches(0.5), Inches(4.3), Inches(12), Inches(0.7),
         size=17, bold=False, color=SLATE, align=PP_ALIGN.CENTER, italic=True)

# Add TechUnited branding bottom
add_text(s, "TECHUNITED:NJ  ·  UPLIFT SUMMER 2026",
         Inches(0.5), Inches(6.8), Inches(12), Inches(0.45),
         size=11, bold=True, color=RGBColor(0x5C, 0x4E, 0xB5), align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
output = "/Users/techunited/Desktop/Uplift-2026-Onboarding-Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
